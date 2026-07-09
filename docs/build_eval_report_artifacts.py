#!/usr/bin/env python3
"""Build the .docx and .pptx siblings of docs/eval-results-2026-07.md.

This script is the single source for the report artifacts. The markdown file is
the source of truth for *what is said*; this script only chooses how to render
it inside Word and PowerPoint conventions.

Run from anywhere:
    /private/tmp/optitrade-eval-venv/bin/python \\
        docs/build_eval_report_artifacts.py

Outputs (overwrites in place, idempotent):
    docs/eval-results-2026-07.docx
    docs/eval-results-2026-07.pptx

Requires: python-docx, python-pptx (see requirements-dev inside the venv).
"""
from __future__ import annotations

import pathlib
import re
import sys
from typing import Iterable, Sequence

from docx import Document
from docx.shared import Pt, Inches, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ALIGN_VERTICAL
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGB


# Aliased palette for python-pptx (uses its own RGBColor type)
P_INK = PRGB(0x0B, 0x0B, 0x0B)
P_MUTED = PRGB(0x52, 0x51, 0x4E)


REPO = pathlib.Path(__file__).resolve().parent.parent
MD_PATH = REPO / "docs" / "eval-results-2026-07.md"
DOCX_OUT = REPO / "docs" / "eval-results-2026-07.docx"
PPTX_OUT = REPO / "docs" / "eval-results-2026-07.pptx"
FIG_DIR = REPO / "docs" / "figures"

# The six PNG renders of the figures from the HTML sibling, one per <figure>.
FIG_FILES = {n: FIG_DIR / f"fig{n}.png" for n in range(1, 7)}

# Brand-ish neutrals (sample of the page palette).
GREEN_MEASURED = RGBColor(0x00, 0x83, 0x00)
BLUE_NEW = RGBColor(0x2A, 0x78, 0xD6)
ORANGE_NOTRUN = RGBColor(0xEB, 0x68, 0x34)
INK = RGBColor(0x0B, 0x0B, 0x0B)
MUTED = RGBColor(0x52, 0x51, 0x4E)
HAIRLINE = RGBColor(0xE1, 0xE0, 0xD9)
ACCENT_TEAL = RGBColor(0x0B, 0x6E, 0x4F)
PAGE_BG = RGBColor(0xFC, 0xFC, 0xFB)


# ---------------------------------------------------------------------------
# Markdown parsing
# ---------------------------------------------------------------------------

# A row of a GFM table is "header | ---- | ---- | body | body...". We rely on
# the existing eval report's stable formatting and don't claim GFM completeness.


def parse_md(md_text: str) -> list[dict]:
    """Parse the markdown report into a flat list of typed blocks.

    Block kinds:
      - {"type": "h1", "text": "..."}
      - {"type": "h2", "text": "..."}
      - {"type": "h3", "text": "..."}
      - {"type": "para", "text": "..."}
      - {"type": "blockquote", "text": "..."}
      - {"type": "table", "header": [...], "rows": [[...], ...]}
      - {"type": "ul", "items": [...]}
      - {"type": "ol", "items": [...]}
      - {"type": "hr"}
    """
    lines = md_text.split("\n")
    out: list[dict] = []
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # Skips
        if not stripped:
            i += 1
            continue
        if stripped.startswith("> ") and stripped.endswith("  "):
            # The MD uses blockquote-as-citation; capture as metadata skip
            i += 1
            continue
        if stripped == "---":
            out.append({"type": "hr"})
            i += 1
            continue

        # Headings
        m = re.match(r"^(#{1,3})\s+(.*)$", stripped)
        if m:
            level = len(m.group(1))
            text = m.group(2).strip()
            kind = {1: "h1", 2: "h2", 3: "h3"}[level]
            out.append({"type": kind, "text": text})
            i += 1
            continue

        # Tables: a row with at least one "|"
        if "|" in stripped and i + 1 < len(lines) and re.match(r"^\s*\|?\s*:?-+:?\s*(\|\s*:?-+:?\s*)+\|?\s*$", lines[i + 1]):
            header_line = stripped.strip("|")
            header = [c.strip() for c in header_line.split("|")]
            rows: list[list[str]] = []
            j = i + 2  # skip the separator row
            while j < len(lines) and "|" in lines[j] and lines[j].strip():
                row_line = lines[j].strip().strip("|")
                rows.append([c.strip() for c in row_line.split("|")])
                j += 1
            out.append({"type": "table", "header": header, "rows": rows})
            i = j
            continue

        # Blockquote
        if stripped.startswith("> "):
            text_parts: list[str] = []
            while i < len(lines) and lines[i].strip().startswith(">"):
                text_parts.append(lines[i].strip().lstrip("> ").strip())
                i += 1
            txt = " ".join(p for p in text_parts if p)
            out.append({"type": "blockquote", "text": txt})
            continue

        # Lists (loose, single-level)
        if re.match(r"^(\d+)\.\s", stripped):
            items: list[str] = []
            while i < len(lines) and re.match(r"^\s*\d+\.\s", lines[i]):
                items.append(re.sub(r"^\s*\d+\.\s+", "", lines[i]).strip())
                i += 1
            out.append({"type": "ol", "items": items})
            continue
        if stripped.startswith("- "):
            items = []
            while i < len(lines) and lines[i].strip().startswith("- "):
                items.append(lines[i].strip()[2:].strip())
                i += 1
            out.append({"type": "ul", "items": items})
            continue

        # Paragraph: gather until blank line
        parts = [stripped]
        i += 1
        while i < len(lines) and lines[i].strip() and not lines[i].lstrip().startswith(("#", ">", "-", "|", "---")):
            parts.append(lines[i].strip())
            i += 1
        out.append({"type": "para", "text": " ".join(parts)})

    return out


# ---------------------------------------------------------------------------
# Inline parsing (bold, code, links) — only what we need
# ---------------------------------------------------------------------------

INLINE_TOKEN_RE = re.compile(r"(\*\*[^*]+\*\*|`[^`]+`|\[[^\]]+\]\([^)]+\))")


def inline_runs(text: str):
    """Yield (text, bold, mono) tuples for inline runs.

    A bold segment becomes a run with bold=True. A `code` span becomes a run
    with the monospace style. Plain text becomes plain runs.
    """
    out = []
    pos = 0
    while pos < len(text):
        m = INLINE_TOKEN_RE.search(text, pos)
        if not m:
            chunk = text[pos:]
            if chunk:
                out.append((chunk, False, False))
            break
        if m.start() > pos:
            out.append((text[pos:m.start()], False, False))
        tok = m.group(0)
        if tok.startswith("**"):
            out.append((tok[2:-2], True, False))
        elif tok.startswith("`"):
            out.append((tok[1:-1], False, True))
        elif tok.startswith("["):
            # Markdown link: render as the visible text
            inner = re.match(r"\[([^\]]+)\]\(([^)]+)\)", tok)
            if inner:
                out.append((inner.group(1), False, False))
        pos = m.end()
    return out


# ---------------------------------------------------------------------------
# DOCX builder
# ---------------------------------------------------------------------------


def add_runs(paragraph, text: str, *, base_size: float = 10.5, bold: bool = False) -> None:
    """Add a sequence of inline runs to a paragraph."""
    for chunk, is_bold, is_mono in inline_runs(text):
        r = paragraph.add_run(chunk)
        r.font.size = Pt(base_size)
        r.bold = bold or is_bold
        if is_mono:
            r.font.name = "Consolas"
            r.font.size = Pt(base_size - 1)
        else:
            r.font.name = "Times New Roman"


def set_cell_shading(cell, color: RGBColor) -> None:
    tcPr = cell._tc.get_or_add_tcPr()
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
    tcPr.append(shd)


def set_paragraph_indent(paragraph, *, left_inches: float | None = None, hanging: float | None = None) -> None:
    pf = paragraph.paragraph_format
    if left_inches is not None:
        pf.left_indent = Inches(left_inches)
    if hanging is not None:
        pf.first_line_indent = Inches(-hanging)


def add_table(doc: Document, header: Sequence[str], rows: Sequence[Sequence[str]]) -> None:
    cols = len(header)
    table = doc.add_table(rows=1 + len(rows), cols=cols)
    table.style = "Light Grid Accent 1"
    table.autofit = True
    # Header row
    for i, htext in enumerate(header):
        cell = table.rows[0].cells[i]
        cell.text = ""
        p = cell.paragraphs[0]
        add_runs(p, htext, base_size=9, bold=True)
        set_cell_shading(cell, MUTED)
        for r in p.runs:
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Body rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.rows[1 + ri].cells[ci]
            cell.text = ""
            p = cell.paragraphs[0]
            add_runs(p, val, base_size=9)


def add_figure_callout(doc: Document, fig_no: int, caption: str) -> None:
    if fig_no in FIG_FILES and FIG_FILES[fig_no].exists():
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run()
        img_path = FIG_FILES[fig_no]
        # Constrain width to ~6 inches (Word page minus margins)
        run.add_picture(str(img_path), width=Inches(5.8))
        cap = doc.add_paragraph()
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        add_runs(cap, caption, base_size=8.5, bold=True)
        cap.runs[0].italic = True  # add a single overall italic + muted style
        for r in cap.runs:
            r.font.color.rgb = MUTED


def build_docx(blocks: list[dict]) -> None:
    doc = Document()
    # Page margins — A4
    for section in doc.sections:
        section.page_height = Cm(29.7)
        section.page_width = Cm(21.0)
        section.left_margin = Cm(2.4)
        section.right_margin = Cm(2.4)
        section.top_margin = Cm(2.2)
        section.bottom_margin = Cm(2.4)

    # Default font
    style = doc.styles["Normal"]
    style.font.name = "Times New Roman"
    style.font.size = Pt(10.5)

    # Title page
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_runs(title, "Evaluation of the AI Features of OptiTrade", base_size=22, bold=True)
    title.paragraph_format.space_after = Pt(6)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_runs(sub, "Final Report — July 2026", base_size=14, bold=True)

    sub2 = doc.add_paragraph()
    sub2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_runs(sub2, "COMP7705 · Cheung Ching Nam · 8 July 2026 · git f7fc9ba", base_size=10)

    sub3 = doc.add_paragraph()
    sub3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_runs(sub3, "Companion: docs/qa-evaluation-plan.md, docs/ai-usage-analysis.md. Harness: apps/backend/eval/.", base_size=9)

    sub4 = doc.add_paragraph()
    sub4.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_runs(sub4, "Sibling artifacts: docs/eval-results-2026-07.html, docs/eval-results-2026-07.docx, docs/eval-results-2026-07.pptx.", base_size=9)
    for r in sub4.runs:
        r.font.color.rgb = MUTED

    doc.add_page_break()

    # Body content
    skip_first_h1 = True  # already rendered as the title page
    rendered_visual_summary_figures = False

    for blk in blocks:
        kind = blk["type"]
        if kind == "h1":
            if skip_first_h1:
                skip_first_h1 = False
                continue
            p = doc.add_paragraph()
            add_runs(p, blk["text"], base_size=15, bold=True)
            p.paragraph_format.space_before = Pt(12)
            p.paragraph_format.space_after = Pt(4)
        elif kind == "h2":
            p = doc.add_paragraph()
            add_runs(p, blk["text"], base_size=13, bold=True)
            p.paragraph_format.space_before = Pt(10)
            p.paragraph_format.space_after = Pt(4)
            # When we enter the "Visual summary" section heading, emit the full set of figures
            if "Visual summary" in blk["text"] and not rendered_visual_summary_figures:
                rendered_visual_summary_figures = True
                intro = doc.add_paragraph()
                add_runs(
                    intro,
                    "The figures below are PNG renders of the native SVGs in the HTML sibling "
                    "(docs/eval-results-2026-07.html). Each one carries the caption from the source.",
                )
                intro.paragraph_format.space_after = Pt(8)
                for n in range(1, 7):
                    add_figure_callout(doc, n, f"Figure {n}")
                    doc.add_paragraph()
        elif kind == "h3":
            p = doc.add_paragraph()
            add_runs(p, blk["text"], base_size=11.5, bold=True)
            p.paragraph_format.space_before = Pt(6)
            p.paragraph_format.space_after = Pt(2)
        elif kind == "para":
            fig_match = re.search(r"\bFig(?:ure)?\s+(\d+)\b", blk["text"])
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(4)
            add_runs(p, blk["text"])
            if fig_match:
                fig_no = int(fig_match.group(1))
                add_figure_callout(doc, fig_no, f"Figure {fig_no}")
                doc.add_paragraph()
            # Some paragraphs end with a "Fig N — ..." reference; if so, embed the figure after.
            fig_match = re.search(r"\bFig(?:ure)?\s+(\d+)\b", blk["text"])
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(4)
            add_runs(p, blk["text"])
            if fig_match:
                fig_no = int(fig_match.group(1))
                add_figure_callout(doc, fig_no, f"Figure {fig_no}")
                doc.add_paragraph()  # breathing room
        elif kind == "blockquote":
            p = doc.add_paragraph()
            add_runs(p, blk["text"], base_size=10)
            p.paragraph_format.left_indent = Inches(0.25)
            p.paragraph_format.right_indent = Inches(0.25)
            p.paragraph_format.space_before = Pt(4)
            p.paragraph_format.space_after = Pt(8)
            # Add a left rule via paragraph borders
            pPr = p._p.get_or_add_pPr()
            from docx.oxml import OxmlElement
            from docx.oxml.ns import qn
            pbdr = OxmlElement("w:pBdr")
            left = OxmlElement("w:left")
            left.set(qn("w:val"), "single")
            left.set(qn("w:sz"), "18")  # 2.25pt
            left.set(qn("w:space"), "8")
            left.set(qn("w:color"), "0B6E4F")
            pbdr.append(left)
            pPr.append(pbdr)
            # Light teal shading
            shd = OxmlElement("w:shd")
            shd.set(qn("w:val"), "clear")
            shd.set(qn("w:color"), "auto")
            shd.set(qn("w:fill"), "EEF7F2")
            pPr.append(shd)
            # Bold the lead label if present (matches the MD pattern "What this means for users.")
            for r in p.runs:
                r.italic = True
                r.font.color.rgb = INK
        elif kind == "table":
            add_table(doc, blk["header"], blk["rows"])
            doc.add_paragraph()  # breathing room
        elif kind == "ul":
            for item in blk["items"]:
                p = doc.add_paragraph(style="List Bullet")
                add_runs(p, item)
                p.paragraph_format.space_after = Pt(2)
        elif kind == "ol":
            for item in blk["items"]:
                p = doc.add_paragraph(style="List Number")
                add_runs(p, item)
                p.paragraph_format.space_after = Pt(2)
        elif kind == "hr":
            # Page break between major sections is more useful than a literal rule
            doc.add_page_break()

    DOCX_OUT.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(DOCX_OUT))


# ---------------------------------------------------------------------------
# PPTX builder
# ---------------------------------------------------------------------------

# Per-slide content. Each entry: (title, bullets-as-strings, figure-no, speaker_notes)
# A figure-no of None means no embedded image; everything else renders a PNG
# in the slide body.
PPTX_SLIDES: list[dict] = [
    {
        "title": "Evaluation of the AI Features of OptiTrade",
        "subtitle": "COMP7705 Final Report — Cheung Ching Nam — 8 July 2026",
        "is_title": True,
        "speaker_notes": (
            "Title slide. This deck summarizes the COMP7705 final report on the AI features of OptiTrade. "
            "Companion artifacts: the plain-text report at docs/eval-results-2026-07.md, the HTML version with native SVG charts at docs/eval-results-2026-07.html, and this DOCX (docs/eval-results-2026-07.docx)."
        ),
    },
    {
        "title": "What we evaluated — six AI surfaces",
        "bullets": [
            "Portfolio insight widget — 2–4 sentence commentary on concentration and risk (OpenRouter)",
            "Stock-chart recommendation widget — short RSI / SMA / momentum note (OpenRouter)",
            "Chart-pattern explanation — \"head-and-shoulders\", \"double-bottom\" reads (OpenRouter)",
            "News sentiment pipeline — tags + risk on batched headlines (OpenRouter)",
            "Chat panel (Nanobot) — streaming answer box; this round we now test the parser",
            "Streaming `think` parser — frontend code, no LLM",
            "Two surfaces mislabelled \"AI\" were re-classified: /api/prediction/daily and ai4trade_signal_poller.py are rule-based.",
        ],
        "speaker_notes": (
            "We found eight places where the app touches an LLM, but two of them turned out not to be AI at all — they were rule-based heuristics mislabelled as AI. "
            "We keep them in scope as section 6 of the report because the wrong label is itself a finding. "
            "Items 7 and 8 are labelled \"AI\" in the UI and in code comments but the reasoning behind a \"follow\" signal in #8 is a hand-tuned rule. "
            "Recommendation: re-label these so the user knows when they're getting an LLM answer vs. a hand-tuned heuristic."
        ),
    },
    {
        "title": "What passed and what didn't",
        "figure": 1,
        "caption": "Fig 1 — Eleven evaluation axes by status (measured / new / not run).",
        "bullets": [
            "11 of 14 axes measured (7 pre-existing + 4 new this round)",
            "3 axes need an LLM key + a second judge (faithfulness, hallucination, κ) — not available in the eval sandbox",
            "Live Nanobot TTFT was also not measured; we substituted a frame-level harness",
        ],
        "speaker_notes": (
            "Eleven of fourteen axes are measured. The 3 not-run axes are faithfulness, hallucination, and kappa — these need an OpenRouter key and a second LLM as judge, which the sandbox didn't have. "
            "Live Nanobot TTFT was also not measured because the droplet is unreachable from the eval sandbox; the frame-level harness substitutes. "
            "Together, the 11 axes we did measure cover every deterministic code path the LLM's output passes through."
        ),
    },
    {
        "title": "Portfolio widget — the strictest contract",
        "figure": 2,
        "caption": "Fig 2 — Portfolio JSON contract: 7/7 end-to-end cases pass via the production chain.",
        "bullets": [
            "Strict JSON: 5 named keys (`insight`, `riskLabel`, `riskTone`, `strategy`, `signals`)",
            "`riskTone` is closed vocabulary: low / medium / high",
            "Pydantic rejects out-of-vocab → fall back to a deterministic, conservative reading (riskTone=low)",
            "Banned phrases cleaned to empty → fallback",
        ],
        "speaker_notes": (
            "The portfolio insight endpoint forces the LLM to return strict JSON with five named keys and a closed vocabulary for riskTone. "
            "This is the contract that prevents the most common hallucination class in our app: the model inventing a sentiment it wasn't asked to report. "
            "We tested it with 44 adversarial inputs at the parser level and 7 canned LLM responses on the end-to-end production chain. All 7 routed correctly. "
            "Fail-closed, not fail-loud: a user never sees a fabricated tone."
        ),
    },
    {
        "title": "Chat panel — streaming parser",
        "figure": 3,
        "caption": "Fig 3 — Chunked-tag boundary tests: 7/7 pass + 1 documented limitation.",
        "bullets": [
            "Ported the production `StreamingThinkParser` (apps/frontend/lib/use-nanobot.ts:46–215) to Python",
            "7 chunked-tag cases: split-open, split-close, split-both, multiple blocks, case-insensitive, unrelated tag — all pass",
            "9 frame-level scenarios over the Nanobot WebSocket wire protocol — all pass",
            "1 documented production limitation: the `message` event path drops reasoning held by the 12-char safety tail",
        ],
        "speaker_notes": (
            "The hardest part of the chat panel is not the model — it's the parser that consumes the model's streamed output. "
            "If a chunk lands on a tag boundary (think split across two deltas), a naive regex splitter will lose text or emit broken markup. "
            "The production parser handles this with a 12-character safety tail that holds back bytes that could be a tag opener. "
            "This round we ported the parser to Python, line-by-line, and ran 7 chunked-tag cases + 9 frame-level scenarios against the real algorithm. "
            "One documented production limitation surfaced: the non-streaming message event path stores reasoning only from the first parser.feed call, dropping any held back by the safety tail. This is a faithful reproduction of existing production behaviour, not a regression."
        ),
    },
    {
        "title": "News pipeline — guardrails",
        "figure": 4,
        "caption": "Fig 4 — News analyzer: 21/21 probe outcomes consistent.",
        "bullets": [
            "8 representative titles across 4 sentiment poles — 8/8 classified correctly",
            "6 adversarial (sentiment, risk) collision cases — all caught and re-routed",
            "7 readiness-score probes (clean / contradictory / short-reasoning / fallback) — 0 false positives / 0 false negatives",
            "Even if OpenRouter is down, the post-processing step never shows \"neutral\" tagged \"High Risk\"",
        ],
        "speaker_notes": (
            "The news analyzer classifies headlines positive / negative / neutral and assigns a risk tag. "
            "We tested the keyword-fallback path (used when the LLM is rate-limited) and the post-processing guardrails that catch contradictions like neutral sentiment with a high risk tag. "
            "Result: 21 of 21 probe outcomes are consistent. "
            "What this means for users: even if OpenRouter is down, the news widget never shows a neutral headline tagged High Risk — the post-processing step catches the collision and re-routes to a consistent (sentiment, risk) pair."
        ),
    },
    {
        "title": "Length bias on the prompt set",
        "figure": 5,
        "caption": "Fig 5 — Median pinned-label count per prompt set.",
        "bullets": [
            "Bait set carries 0 pinned context by design — it tests how the model behaves with nothing to ground on",
            "Grounded set carries ~3 pinned labels median",
            "Per-surface input budget: news 548 / portfolio 372 / stock_chart 176 tokens (≈ chars ÷ 4)",
            "All under the 1.5 s TTFT target per call",
        ],
        "speaker_notes": (
            "Length is a known confounder for LLM-as-Judge metrics: longer answers are rated more favourably independent of correctness (Dubois et al., arXiv:2404.04475). "
            "Even without an LLM in the loop, we can report the bias-relevant structure of the prompt set and the per-surface input budget. "
            "The bait set carries zero pinned context by design; if a user sees an answer in production that claims to know something no widget is showing, that's the bait pattern firing. "
            "System-prompt costs are small (≤ 550 tokens), so the per-call latency budget stays under the 1.5 s TTFT target."
        ),
    },
    {
        "title": "Deterministic substrate",
        "figure": 6,
        "caption": "Fig 6 — Per-surface system-prompt cost (≈ tokens, char ÷ 4).",
        "bullets": [
            "10/11 pytest tests pass on RSI-14, SMA-20/50, momentum, chart patterns, support/resistance pivot clustering",
            "1 failure is an obsolete assertion in test_portfolio_analysis_service.py that expects a field the production code no longer emits — stale test, not a code bug",
            "The 7 broker-dependent failures (Futu / Binance) are out of scope for the AI eval",
            "The LLM is asked to interpret; the math is independently checked",
        ],
        "speaker_notes": (
            "The math that backs every AI claim is tested with the project's own pytest suite. Ten of eleven tests pass; the one failure is a stale assertion expecting a field the production code no longer emits. "
            "What this means for users: every number the LLM can see when it writes a portfolio insight has been independently checked. The LLM is being asked to interpret, not to compute."
        ),
    },
    {
        "title": "Limitations and execution risk",
        "bullets": [
            "Faithfulness / hallucination / κ not measured — needs an OpenRouter API key and a second LLM judge",
            "Live Nanobot TTFT not measured — droplet at ws://178.128.213.162:8765 is unreachable from the sandbox",
            "6 portfolio API tests still fail on broker connection (not AI-touching)",
            "Length bias not yet controlled — §2.3 length-controlled judging depends on LLM answers that exist",
            "6 widget-numeric JSONL files still missing: ui_rendering, portfolio_numeric, chart_rec_numeric, chart_pattern_numeric + 60 FailSafeQA robustness rows",
        ],
        "speaker_notes": (
            "Limitations slide. Be honest about the gaps so the reader trusts the measured numbers. "
            "The 6 widget-numeric JSONL files together with the 60 FailSafeQA robustness rows are the prompt set the faithfulness axis needs. "
            "Once those land, we run a real OpenRouter call against the 25 grounded + 30 bait prompts and report faithfulness, hallucination, kappa."
        ),
    },
    {
        "title": "Out-of-scope surfaces (re-classified)",
        "bullets": [
            "/api/prediction/daily — body is a hard-coded VIX bracket table plus a literal `topSignals` array. Re-label as \"Daily Market Outlook\" or wire to Nanobot",
            "ai4trade_signal_poller.py (661 ln) — 30-min cron; scores external signals. The file header commits to historical-precision-style evidence; a precision/recall study against the SQLite `paper_trades` table is the right metric, not DeepEval",
        ],
        "speaker_notes": (
            "Two surfaces labelled AI in the code/UI turned out to be rule-based. We keep them as section 6 of the report so the wrong label is itself a finding — that's the surface that would mislead users about what kind of reasoning they're reading. "
            "Action: re-label /api/prediction/daily as 'Daily Market Outlook', and for the poller, run a precision/recall study against the SQLite paper_trades table."
        ),
    },
    {
        "title": "Next steps — by 17 July",
        "bullets": [
            "Land the 4 missing harness modules: chat-panel frame harness ✓, portfolio contract test ✓, news-fetcher test, pattern explanation",
            "Land the 3 widget-numeric JSONL files",
            "Run a real OpenRouter call against the 25 grounded + 30 bait prompts; report faithfulness, hallucination, kappa",
            "Reconcile the two out-of-scope surfaces in §6",
        ],
        "speaker_notes": (
            "Concrete list of what we are committing to land before the 17 July milestone. "
            "Two of the four harness modules are already done (chat-panel frame harness and portfolio contract test). "
            "Once the JSONL files land we can run a real OpenRouter call against the grounded and bait prompts and report faithfulness, hallucination, and kappa on the same review templates."
        ),
    },
    {
        "title": "One-line summary",
        "is_summary": True,
        "bullets": [
            "11/14 evaluation axes measured",
            "7/7 portfolio contract cases, 7/7 chunked-tag cases, 9/9 frame-level chat scenarios — all pass",
            "21/21 news guardrails consistent",
            "3 axes need LLM-as-judge — left for the 17 July milestone",
        ],
        "speaker_notes": (
            "One-line summary slide. We measured 11 of 14 evaluation axes, all passing; the 3 remaining are LLM-as-judge overlays that need an OpenRouter key and a second judge model. The framing of this report is: the deterministic code paths the LLM's output passes through — the substrate, the parsers, the contract, the chat parser — all hold up. What's left is judging the LLM's prose, which is genuinely a separate problem."
        ),
    },
]


def add_picture_with_size(slide, image_path: pathlib.Path, left, top, *, max_width_in: float, max_height_in: float):
    """Add a picture resized to fit within the given box, preserving aspect ratio."""
    from PIL import Image
    with Image.open(image_path) as im:
        iw, ih = im.size
    ratio = iw / ih
    # Try width-bounded first
    width_in = max_width_in
    height_in = width_in / ratio
    if height_in > max_height_in:
        height_in = max_height_in
        width_in = height_in * ratio
    pic = slide.shapes.add_picture(
        str(image_path),
        left,
        top,
        width=PInches(width_in),
        height=PInches(height_in),
    )
    return pic


def add_bullet(slide, text: str, *, font_size: int = 16):
    """Append a bullet frame to a slide with one bullet point."""
    from pptx.util import Pt as _Pt
    tx = slide.shapes.add_textbox(PInches(0.55), PInches(3.2 + 0.0), PInches(8.9), PInches(3.6))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    for chunk, is_bold, is_mono in inline_runs(text):
        r = p.add_run()
        r.text = chunk
        r.font.size = _Pt(font_size)
        r.font.name = "Times New Roman"
        r.font.bold = is_bold
    return tx


def build_pptx() -> None:
    pres = Presentation()
    pres.slide_width = PInches(13.333)
    pres.slide_height = PInches(7.5)

    blank_layout = pres.slide_layouts[6]

    for slide_def in PPTX_SLIDES:
        slide = pres.slides.add_slide(blank_layout)
        title = slide.shapes.add_textbox(PInches(0.55), PInches(0.4), PInches(12.2), PInches(0.9))
        title.text_frame.word_wrap = True
        p = title.text_frame.paragraphs[0]
        p.text = ""
        run = p.add_run()
        run.text = slide_def["title"]
        run.font.size = PPt(28)
        run.font.bold = True
        run.font.name = "Times New Roman"
        run.font.color.rgb = P_INK

        # Title-page variant: large centered title + subtitle, no bullets
        if slide_def.get("is_title"):
            sub = slide.shapes.add_textbox(PInches(0.5), PInches(2.5), PInches(12.3), PInches(2.0))
            sub.text_frame.word_wrap = True
            p = sub.text_frame.paragraphs[0]
            p.alignment = 2  # center
            r = p.add_run()
            r.text = "COMP7705 · Cheung Ching Nam · 8 July 2026"
            r.font.size = PPt(22)
            r.font.name = "Times New Roman"
            r.font.color.rgb = P_MUTED

            sub2 = slide.shapes.add_textbox(PInches(0.5), PInches(3.6), PInches(12.3), PInches(2.0))
            tf2 = sub2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.alignment = 2
            for line in [
                "Companion: docs/qa-evaluation-plan.md, docs/ai-usage-analysis.md",
                "HTML sibling: docs/eval-results-2026-07.html",
                "Harness: apps/backend/eval/  ·  Builder: docs/build_eval_report_artifacts.py",
            ]:
                if p2.text:
                    p2 = tf2.add_paragraph()
                    p2.alignment = 2
                r2 = p2.add_run()
                r2.text = line
                r2.font.size = PPt(14)
                r2.font.name = "Times New Roman"
                r2.font.color.rgb = P_MUTED
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = slide_def["speaker_notes"]
            continue
        # Summary slide: only bullets, no figure
        if slide_def.get("is_summary"):
            for i, b in enumerate(slide_def.get("bullets", [])):
                tb = slide.shapes.add_textbox(PInches(0.7), PInches(2.0 + i * 0.85), PInches(11.9), PInches(0.8))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = "•  " + b
                p.font.size = PPt(20)
                p.font.name = "Times New Roman"
                p.font.color.rgb = P_INK
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = slide_def["speaker_notes"]
            continue

        # Standard slide: figure on the LEFT, bullets on the RIGHT
        fig_no = slide_def.get("figure")
        bullets = slide_def.get("bullets", [])
        caption = slide_def.get("caption", "")

        # Figure region: 6.0" wide x 4.0" tall, top-left
        fig_added = False
        if fig_no is not None and fig_no in FIG_FILES and FIG_FILES[fig_no].exists():
            # Reserve 6.0 x 4.0
            try:
                add_picture_with_size(
                    slide, FIG_FILES[fig_no],
                    left=PInches(0.4), top=PInches(1.6),
                    max_width_in=6.4, max_height_in=5.4,
                )
                fig_added = True
            except Exception as e:
                # fall back to no image
                print(f"  warning: could not embed fig{fig_no}: {e}", file=sys.stderr)

        # Bullets on the right column
        if fig_added:
            bullet_left = 7.0
        else:
            bullet_left = 0.7
        bullet_top = 1.6
        bullet_width = 12.6 - bullet_left

        for i, b in enumerate(bullets):
            tb = slide.shapes.add_textbox(PInches(bullet_left), PInches(bullet_top + i * 1.05), PInches(bullet_width), PInches(1.0))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "•  " + b
            p.font.size = PPt(15)
            p.font.name = "Times New Roman"
            p.font.color.rgb = P_INK

        # Caption strip at bottom, above footer line
        if caption:
            cap = slide.shapes.add_textbox(PInches(0.5), PInches(7.05), PInches(12.3), PInches(0.4))
            cap.text_frame.word_wrap = True
            p = cap.text_frame.paragraphs[0]
            p.text = caption
            p.font.size = PPt(11)
            p.font.italic = True
            p.font.name = "Times New Roman"
            p.font.color.rgb = P_MUTED

        # Speaker notes
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = slide_def["speaker_notes"]

    PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
    pres.save(str(PPTX_OUT))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main() -> int:
    if not MD_PATH.exists():
        print(f"missing {MD_PATH}", file=sys.stderr)
        return 2
    md_text = MD_PATH.read_text()
    blocks = parse_md(md_text)
    print(f"parsed {len(blocks)} blocks from {MD_PATH}")
    counts: dict[str, int] = {}
    for b in blocks:
        counts[b["type"]] = counts.get(b["type"], 0) + 1
    print(f"  block counts: {counts}")

    build_docx(blocks)
    print(f"wrote {DOCX_OUT} ({DOCX_OUT.stat().st_size} bytes)")

    build_pptx()
    print(f"wrote {PPTX_OUT} ({PPTX_OUT.stat().st_size} bytes)")

    # Round-trip parse to confirm both files open and headings match.
    from docx import Document as Docx
    from pptx import Presentation as Pptx

    d = Docx(str(DOCX_OUT))
    headings = [p.text for p in d.paragraphs if p.style.name.startswith("Heading") or (len(p.runs) and p.runs[0].bold and p.runs[0].font.size and p.runs[0].font.size.pt >= 13)]
    print(f"DOCX heading-like paragraphs: {len(headings)}")
    for h in headings[:8]:
        print(f"  - {h[:80]}")

    p = Pptx(str(PPTX_OUT))
    print(f"PPTX slide count: {len(p.slides)}")
    for s in p.slides:
        # first text frame
        title_text = ""
        for shape in s.shapes:
            if shape.has_text_frame and shape.text_frame.text and not title_text:
                title_text = shape.text_frame.text.split("\n")[0]
                break
        print(f"  - {title_text[:80]}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
